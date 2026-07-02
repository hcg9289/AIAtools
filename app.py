import os
import io
import json
import base64
import uuid
import re
import threading
import subprocess
import requests
import time
import fitz  # PyMuPDF
from PIL import Image
try:
    import pytesseract
except ImportError:  # pragma: no cover - runtime dependency is installed in Docker
    pytesseract = None
from datetime import datetime, timedelta
from flask import Flask, render_template, request, jsonify, send_file, redirect
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
os.makedirs('uploads', exist_ok=True)
os.makedirs('outputs', exist_ok=True)

app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

VAULT_AI_URL = os.environ.get('VAULT_AI_URL', 'http://wa-vault-1006:5001')
PPT_AI_ENDPOINT = f"{VAULT_AI_URL}/api/v1/ai/ppt/1008"
PPT_GENERATE_ENDPOINT = f"{VAULT_AI_URL}/api/v1/ai/ppt/generate"
VAULT_AUTH_URL = os.environ.get('VAULT_AUTH_URL', 'http://wa-vault-1006:5001/api/v1/token/validate')
SESSION_TTL_SECONDS = int(os.environ.get('SESSION_TTL_SECONDS', str(20 * 60)))
AUTH_SESSIONS = {}  # {sid: {"uid": str, "expiry": datetime, "ott": str}}
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'pdf', 'ppt', 'pptx'}
PREVIEW_BOTS = ['WhatsApp', 'facebookexternalhit', 'Twitterbot', 'LinkedInBot', 'Slackbot']
VAULT_AUTH_TIMEOUT_SECONDS = 15
VAULT_AUTH_RETRY_ATTEMPTS = 2
VAULT_AUTH_RETRY_BACKOFF_SECONDS = 0.4
AUTH_RETRY_PAGE_SECONDS = 2

# Async task store — {task_id: {"status": "processing"|"done"|"error", "result": dict, "created": datetime}}
TASKS = {}
TASKS_LOCK = threading.Lock()

TASK_TTL_SECONDS = 1800  # 30 minutes
PPT_FONT_FAMILY = os.environ.get('PPT_FONT_FAMILY', 'Noto Sans CJK TC')
PPT_FALLBACK_FONT_FAMILY = os.environ.get('PPT_FALLBACK_FONT_FAMILY', 'Microsoft JhengHei')
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
GF_TEMPLATE_PATH = os.environ.get(
    'GF_TEMPLATE_PATH',
    os.path.join(BASE_DIR, 'assets', 'gf', 'gf_template.pptx')
)
GF_MODEL_PAGE_PATH = os.environ.get(
    'GF_MODEL_PAGE_PATH',
    os.path.join(BASE_DIR, 'research', 'gf_withdrawal_model.html')
)
GF_WITHDRAWAL_TARGET_YEAR = 20
GF_OCR_ZOOM = float(os.environ.get('GF_OCR_ZOOM', '4.0'))
GF_OCR_LANG = os.environ.get('GF_OCR_LANG', 'eng')
GF_OCR_MIN_CONFIDENCE = float(os.environ.get('GF_OCR_MIN_CONFIDENCE', '45'))

GF_CROP_RULES = {
    'proposal_summary': {
        'terms': ('建議書摘要', '保障摘要'),
        'crop': (0.00, 0.265, 1.00, 0.735),
    },
    'summary_table': {
        'terms': ('基本計劃', '說明摘要'),
        'crop': (0.00, 0.03, 1.00, 0.78),
    },
    'withdrawal_amount_table': {
        'terms': ('細分之保證及非保證現金提取金額',),
        'crop': (0.06, 0.18, 0.94, 0.91),
    },
    'withdrawal_surrender_table': {
        'terms': ('現金提取後之退保發還金額',),
        'crop': (0.035, 0.13, 0.96, 0.92),
    },
}

LAYOUT_BOXES = {
    'left': (0.72, 0.70, 5.45, 5.92),
    'right': (7.18, 0.70, 5.45, 5.92),
    'center': (1.38, 1.04, 10.56, 5.38),
    'top': (0.90, 0.60, 11.55, 3.08),
    'bottom': (0.90, 3.38, 11.55, 3.30),
    'full': (0.78, 0.70, 11.78, 5.92),
}


def _cleanup_old_tasks():
    now = datetime.now()
    with TASKS_LOCK:
        stale = [k for k, v in TASKS.items() if (now - v['created']).total_seconds() > TASK_TTL_SECONDS]
        for k in stale:
            del TASKS[k]


def get_client_ip():
    cf_ip = request.headers.get('cf-connecting-ip')
    if cf_ip:
        return cf_ip.strip()
    forwarded = request.headers.get('x-forwarded-for')
    if forwarded:
        return forwarded.split(',')[0].strip()
    real_ip = request.headers.get('x-real-ip')
    if real_ip:
        return real_ip.strip()
    return request.remote_addr or '127.0.0.1'


def expects_json_response():
    accept = request.headers.get('Accept', '').lower()
    return request.path.startswith('/api/') or request.is_json or 'application/json' in accept or request.method != 'GET'


def auth_forbidden_response(message='Unauthorized'):
    if expects_json_response():
        return jsonify({'success': False, 'error': message}), 403
    return (
        '<!DOCTYPE html><html lang="zh-Hant"><head><meta charset="utf-8"><title>授權失敗</title></head>'
        '<body style="font-family:sans-serif;text-align:center;padding-top:50px">'
        f'<h2>授權失敗 (1008)</h2><p>{message}</p><p>請回 WhatsApp 重新取得工具連結。</p>'
        '</body></html>',
        403,
        {'Content-Type': 'text/html; charset=utf-8'}
    )


def auth_retry_response(status_code, message):
    if expects_json_response():
        return jsonify({'success': False, 'error': {'code': status_code, 'message': message}}), 503
    return (
        '<!DOCTYPE html><html lang="zh-Hant"><head><meta charset="utf-8">'
        f'<meta http-equiv="refresh" content="{AUTH_RETRY_PAGE_SECONDS}">'
        '<title>OTT 驗證中</title></head>'
        '<body style="font-family:sans-serif;text-align:center;padding-top:50px">'
        '<h2>OTT 驗證暫時未確認 (1008)</h2>'
        f'<p>{message}</p><p>系統會自動重試，或你可以重新整理同一條連結。</p>'
        '</body></html>',
        503,
        {
            'Content-Type': 'text/html; charset=utf-8',
            'Retry-After': str(AUTH_RETRY_PAGE_SECONDS),
            'Cache-Control': 'no-store, max-age=0',
        }
    )


def is_vault_timeout_error(exc):
    return isinstance(exc, requests.exceptions.Timeout) or 'timeout' in str(exc).lower() or 'timed out' in str(exc).lower()


def validate_vault_ott(ott):
    token_prefix = (ott or '')[:8]
    real_ip = get_client_ip()
    request_id = uuid.uuid4().hex
    saw_timeout = False

    for attempt in range(1, VAULT_AUTH_RETRY_ATTEMPTS + 1):
        try:
            resp = requests.get(
                VAULT_AUTH_URL,
                params={'token': ott, 'request_id': request_id},
                headers={'CF-Connecting-IP': real_ip},
                timeout=VAULT_AUTH_TIMEOUT_SECONDS
            )
            if resp.status_code == 200:
                data = resp.json()
                if data.get('valid'):
                    app.logger.info('OTT validation result=valid token_prefix=%s attempt=%s real_ip=%s vault_code=%s', token_prefix, attempt, real_ip, resp.status_code)
                    return 'valid', data, attempt
                result = 'auth_uncertain' if saw_timeout else 'invalid'
                app.logger.info('OTT validation result=%s token_prefix=%s attempt=%s real_ip=%s vault_code=%s', result, token_prefix, attempt, real_ip, resp.status_code)
                return result, data, attempt
            if resp.status_code in (401, 403):
                result = 'auth_uncertain' if saw_timeout else 'invalid'
                app.logger.info('OTT validation result=%s token_prefix=%s attempt=%s real_ip=%s vault_code=%s', result, token_prefix, attempt, real_ip, resp.status_code)
                return result, None, attempt
            if resp.status_code >= 500 and attempt < VAULT_AUTH_RETRY_ATTEMPTS:
                app.logger.info('OTT validation result=temporary_error_retry token_prefix=%s attempt=%s real_ip=%s vault_code=%s', token_prefix, attempt, real_ip, resp.status_code)
                time.sleep(VAULT_AUTH_RETRY_BACKOFF_SECONDS * attempt)
                continue
            app.logger.info('OTT validation result=temporary_error token_prefix=%s attempt=%s real_ip=%s vault_code=%s', token_prefix, attempt, real_ip, resp.status_code)
            return 'temporary_error', None, attempt
        except Exception as exc:
            if is_vault_timeout_error(exc):
                saw_timeout = True
                result = 'timeout_retry' if attempt < VAULT_AUTH_RETRY_ATTEMPTS else 'timeout'
                app.logger.info('OTT validation result=%s token_prefix=%s attempt=%s real_ip=%s exception=%s', result, token_prefix, attempt, real_ip, type(exc).__name__)
                if attempt < VAULT_AUTH_RETRY_ATTEMPTS:
                    time.sleep(VAULT_AUTH_RETRY_BACKOFF_SECONDS * attempt)
                    continue
                return 'timeout', None, attempt
            result = 'auth_uncertain' if saw_timeout else 'temporary_error'
            app.logger.info('OTT validation result=%s token_prefix=%s attempt=%s real_ip=%s exception=%s', result, token_prefix, attempt, real_ip, type(exc).__name__)
            return result, None, attempt
    return ('auth_uncertain' if saw_timeout else 'temporary_error'), None, VAULT_AUTH_RETRY_ATTEMPTS


@app.before_request
def verify_ott_access():
    if request.path.startswith('/static') or request.path in ('/health',):
        return
    if request.path.endswith('.css') or request.path.endswith('.js') or request.path.endswith('.png') or request.path.endswith('.jpg') or request.path.endswith('.ico'):
        return

    ua = request.headers.get('User-Agent', '')
    if any(bot in ua for bot in PREVIEW_BOTS):
        if expects_json_response():
            return jsonify({'success': False, 'error': 'Unauthorized'}), 403
        return (
            '<!DOCTYPE html><html><head><meta charset="utf-8">'
            '<title>PPT Generator</title>'
            '<meta property="og:title" content="PPT Generator"/>'
            '<meta property="og:description" content="請透過 WhatsApp 取得專屬連結後使用。"/>'
            '</head><body></body></html>',
            200,
            {'Content-Type': 'text/html; charset=utf-8'}
        )

    sid = request.cookies.get('auth_sid')
    if sid:
        session = AUTH_SESSIONS.get(sid)
        if session and session['expiry'] > datetime.now():
            return

    ott = request.args.get('ott')
    if ott:
        try:
            status, data, attempt = validate_vault_ott(ott)
            if status == 'valid':
                new_sid = str(uuid.uuid4())
                AUTH_SESSIONS[new_sid] = {
                    'uid': data.get('uid'),
                    'tool': data.get('tool'),
                    'ip': data.get('ip'),
                    'expiry': datetime.now() + timedelta(seconds=SESSION_TTL_SECONDS),
                }
                clean_url = request.path
                out = redirect(clean_url)
                out.set_cookie('auth_sid', new_sid, max_age=SESSION_TTL_SECONDS, httponly=True, samesite='Lax')
                return out
            if status == 'invalid':
                return auth_forbidden_response('OTT 連結已失效，請從 WhatsApp 重新取得工具連結。')
            return auth_retry_response(status, '授權驗證服務暫時未確認，請重試同一條連結。')
        except Exception:
            return auth_retry_response('temporary_error', '授權驗證服務暫時未確認，請稍後重試同一條連結。')

    return auth_forbidden_response('Unauthorized')


# ── AI 背景圖模式映射（1006 新端點用）───────────────────────
ASPECT_RATIO_MAP = {
    '16:9': (16, 9),
    '4:3': (4, 3),
    '1:1': (1, 1),
    '3:4': (3, 4),
    '9:16': (9, 16),
}


def _filename_extension(filename):
    return os.path.splitext(filename or '')[1].lower().lstrip('.')


def allowed_file(filename):
    return _filename_extension(filename) in ALLOWED_EXTENSIONS


def _safe_upload_filename(original_filename):
    ext = _filename_extension(original_filename)
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"不支援的檔案格式：{original_filename}")
    return f"upload_{uuid.uuid4()}.{ext}"


def pdf_to_base64_images(pdf_path):
    images = []
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        buffered = io.BytesIO()
        img.save(buffered, format="JPEG")
        img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
        images.append(f"data:image/jpeg;base64,{img_str}")
    return images


def convert_ppt_to_pdf(ppt_path, output_dir):
    try:
        ppt_path = os.path.abspath(ppt_path)
        output_dir = os.path.abspath(output_dir)
        subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', ppt_path, '--outdir', output_dir],
            check=True, capture_output=True
        )
        filename = os.path.basename(ppt_path)
        pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
        return os.path.join(output_dir, pdf_filename)
    except Exception as e:
        print(f"Error converting PPT to PDF: {e}")
        return None


def _convert_pptx_to_pdf_or_raise(pptx_path):
    output_dir = os.path.dirname(pptx_path)
    pdf_path = convert_ppt_to_pdf(pptx_path, output_dir)
    if not pdf_path or not os.path.exists(pdf_path):
        raise ValueError('PPT 已生成，但 PDF 轉換失敗。請稍後重試或下載 PPTX。')
    return pdf_path


def process_file_to_images(filepath):
    ext = _filename_extension(filepath)
    if ext not in ALLOWED_EXTENSIONS:
        return []
    if ext in ['png', 'jpg', 'jpeg']:
        with open(filepath, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
            mime_ext = 'jpeg' if ext == 'jpg' else ext
            return [f"data:image/{mime_ext};base64,{encoded_string}"]
    elif ext == 'pdf':
        return pdf_to_base64_images(filepath)
    elif ext in ['ppt', 'pptx']:
        pdf_path = convert_ppt_to_pdf(filepath, app.config['UPLOAD_FOLDER'])
        if pdf_path:
            images = pdf_to_base64_images(pdf_path)
            os.remove(pdf_path)
            return images
        return []
    return []


def _safe_hex(hex_code, fallback='#111827'):
    value = str(hex_code or fallback).strip()
    if not value.startswith('#'):
        value = f'#{value}'
    if len(value) == 4:
        value = '#' + ''.join(ch * 2 for ch in value[1:])
    if len(value) != 7:
        return fallback
    try:
        int(value[1:], 16)
    except ValueError:
        return fallback
    return value.upper()


def hex_to_rgb(hex_code):
    hex_code = _safe_hex(hex_code).lstrip('#')
    return tuple(int(hex_code[i:i+2], 16) for i in (0, 2, 4))


def _rgb(hex_code):
    return RGBColor(*hex_to_rgb(hex_code))


def _set_run_font(run, size_pt, color_hex, bold=False, font_name=None):
    font = run.font
    font.name = font_name or PPT_FONT_FAMILY
    font.size = Pt(size_pt)
    font.bold = bold
    font.color.rgb = _rgb(color_hex)

    resolved_font = font_name or PPT_FONT_FAMILY
    r_pr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        font_elem = r_pr.find(qn(tag))
        if font_elem is None:
            font_elem = OxmlElement(tag)
            r_pr.append(font_elem)
        font_elem.set('typeface', resolved_font)


def _add_textbox(slide, box, vertical_anchor=MSO_ANCHOR.TOP):
    x, y, w, h = box
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return shape


def _paragraph(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, space_after=8):
    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = str(text or '').strip()
    _set_run_font(run, size, color, bold=bold)
    return p


def _set_slide_background(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_hex)


def _move_shape_to_back(slide, shape):
    sp_tree = slide.shapes._spTree
    elem = shape._element
    sp_tree.remove(elem)
    sp_tree.insert(2, elem)


def _add_background_image(prs, slide, base64_img_str):
    """Add a full-slide cover image without stretching the source aspect ratio."""
    try:
        raw = _strip_data_uri_prefix(base64_img_str)
        img_bytes = base64.b64decode(raw)
        pil_img = Image.open(io.BytesIO(img_bytes))
        src_w, src_h = pil_img.size
        if src_w <= 0 or src_h <= 0:
            return False

        slide_aspect = float(prs.slide_width) / float(prs.slide_height)
        img_aspect = src_w / src_h

        img_stream = io.BytesIO(img_bytes)
        pic = slide.shapes.add_picture(
            img_stream,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        if img_aspect > slide_aspect:
            crop = max(0, min(0.49, (1 - (slide_aspect / img_aspect)) / 2))
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_aspect < slide_aspect:
            crop = max(0, min(0.49, (1 - (img_aspect / slide_aspect)) / 2))
            pic.crop_top = crop
            pic.crop_bottom = crop

        _move_shape_to_back(slide, pic)
        return True

    except Exception as e:
        print(f"Background image failed: {e}")
        return False


def _coerce_layout(layout_value, index):
    if isinstance(layout_value, dict):
        layout_value = layout_value.get('text_position') or layout_value.get('type') or layout_value.get('name')
    layout = str(layout_value or '').strip().lower().replace('-', '_')
    aliases = {
        'left_text': 'left',
        'right_text': 'right',
        'center_text': 'center',
        'title': 'center',
        'cover': 'center',
        'closing': 'center',
    }
    layout = aliases.get(layout, layout)
    if layout not in LAYOUT_BOXES:
        layout = 'center' if index == 0 else 'left'
    return layout


def _coerce_body_blocks(slide_data):
    blocks = slide_data.get('body_blocks')
    if isinstance(blocks, list) and blocks:
        return blocks

    anchors = slide_data.get('text_anchors')
    if isinstance(anchors, list):
        anchor_blocks = []
        for anchor in anchors:
            if not isinstance(anchor, dict):
                continue
            anchor_type = str(anchor.get('type') or anchor.get('kind') or '').lower()
            if anchor_type == 'title':
                continue
            text = str(anchor.get('text') or anchor.get('value') or '').strip()
            supporting = str(anchor.get('supporting_text') or anchor.get('label') or '').strip()
            if text and supporting:
                anchor_blocks.append({'type': 'bullet', 'text': f'{text}: {supporting}'})
            elif text:
                anchor_blocks.append({'type': 'bullet', 'text': text})
        if anchor_blocks:
            return anchor_blocks

    content = slide_data.get('content', [])
    if isinstance(content, str):
        content = [content]
    if isinstance(content, list):
        return [{'type': 'bullet', 'text': item} for item in content if str(item).strip()]
    return []


def _clamp_float(value, default, min_value, max_value):
    try:
        value = float(value)
    except (TypeError, ValueError):
        value = default
    return max(min_value, min(max_value, value))


def _coerce_anchor_align(value):
    align = str(value or 'left').strip().lower()
    if align in ('center', 'middle'):
        return PP_ALIGN.CENTER
    if align in ('right', 'end'):
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def _coerce_bool(value, default=False):
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    text = str(value).strip().lower()
    if text in ('1', 'true', 'yes', 'y', 'on'):
        return True
    if text in ('0', 'false', 'no', 'n', 'off'):
        return False
    return default


def _coerce_text_anchors(slide_data, defaults):
    anchors = slide_data.get('text_anchors')
    if not isinstance(anchors, list):
        return []

    def anchor_priority(raw):
        if not isinstance(raw, dict):
            return 99
        anchor_type = str(raw.get('type') or raw.get('kind') or '').strip().lower()
        if anchor_type == 'title':
            return 0
        if str(raw.get('evidence_ref') or raw.get('evidence_id') or '').strip():
            return 1
        if str(raw.get('bind_to') or raw.get('visual_element') or '').strip():
            return 2
        return 3

    normalized = []
    for raw in sorted(anchors, key=anchor_priority)[:8]:
        if not isinstance(raw, dict):
            continue
        text = str(raw.get('text') or raw.get('value') or '').strip()
        supporting = str(raw.get('supporting_text') or raw.get('label') or '').strip()
        if not text and not supporting:
            continue

        box = raw.get('box') if isinstance(raw.get('box'), dict) else raw
        x = _clamp_float(box.get('x'), 0.08, 0.02, 0.96)
        y = _clamp_float(box.get('y'), 0.08, 0.02, 0.96)
        w = _clamp_float(box.get('w'), 0.22, 0.05, 0.72)
        h = _clamp_float(box.get('h'), 0.12, 0.04, 0.42)
        w = min(w, 0.98 - x)
        h = min(h, 0.98 - y)
        if w < 0.04 or h < 0.04:
            continue

        anchor_type = str(raw.get('type') or raw.get('kind') or 'label').strip().lower()
        default_color = defaults['title_color'] if anchor_type == 'title' else defaults['content_color']
        font_size = _clamp_float(raw.get('font_size'), 28 if anchor_type == 'title' else 16, 8, 44)
        if anchor_type == 'title':
            w = max(w, min(0.62 if len(text) > 18 else 0.50, 0.98 - x))
            h = max(h, 0.15 if len(text) > 18 else 0.12)
            font_size = min(font_size, 32 if len(text) > 18 else 34)
        bind_to = str(raw.get('bind_to') or raw.get('visual_element') or '').strip()
        connector_default = bool(bind_to and anchor_type != 'title')
        pad_enabled = _coerce_bool(raw.get('pad'), anchor_type != 'title')
        if anchor_type != 'title':
            w = min(w, 0.34)
            h = min(h, 0.20)
            font_size = min(font_size, 18)
        if pad_enabled and (w > 0.34 or h > 0.20 or (w * h) > 0.062):
            pad_enabled = False
        normalized.append({
            'id': str(raw.get('id') or f'anchor_{len(normalized) + 1}'),
            'type': anchor_type,
            'text': text,
            'supporting_text': supporting,
            'bind_to': bind_to,
            'evidence_ref': str(raw.get('evidence_ref') or raw.get('evidence_id') or '').strip(),
            'placement_reason': str(raw.get('placement_reason') or '').strip(),
            'box': (x, y, w, h),
            'align': _coerce_anchor_align(raw.get('align')),
            'font_size': font_size,
            'color': _safe_hex(raw.get('color') or default_color, default_color),
            'supporting_color': _safe_hex(raw.get('supporting_color') or defaults['content_color'], defaults['content_color']),
            'pad': pad_enabled,
            'pad_color': _safe_hex(raw.get('pad_color') or defaults['anchor_pad_color'], defaults['anchor_pad_color']),
            'pad_opacity': _clamp_float(raw.get('pad_opacity'), 0.48, 0.16, 0.68),
            'connector_required': _coerce_bool(raw.get('connector_required'), connector_default),
            'connector': raw.get('connector') if isinstance(raw.get('connector'), dict) else None,
            'target_point': raw.get('target_point') if isinstance(raw.get('target_point'), dict) else None,
        })

    return normalized


def _normalize_slide(slide_data, index, deck_defaults):
    layout = _coerce_layout(slide_data.get('layout') or slide_data.get('layout_type'), index)
    title = str(slide_data.get('title') or f'Slide {index + 1}').strip()
    body_blocks = _coerce_body_blocks(slide_data)

    colors = deck_defaults.get('colors', {})
    bg_color = _safe_hex(slide_data.get('bg_color') or colors.get('background'), '#111827')
    title_color = _safe_hex(slide_data.get('title_color') or colors.get('title'), '#FFFFFF')
    content_color = _safe_hex(slide_data.get('content_color') or colors.get('content'), '#E5E7EB')
    accent_color = _safe_hex(slide_data.get('accent_color') or colors.get('accent'), '#60A5FA')
    anchor_pad_color = _safe_hex(slide_data.get('anchor_pad_color') or colors.get('anchor_pad') or '#FFFFFF', '#FFFFFF')

    overlay = slide_data.get('overlay') if isinstance(slide_data.get('overlay'), dict) else {}
    overlay_color = _safe_hex(overlay.get('color') or slide_data.get('overlay_color') or '#050816', '#050816')
    overlay_opacity = overlay.get('opacity', slide_data.get('overlay_opacity', 0.46))
    try:
        overlay_opacity = float(overlay_opacity)
    except (TypeError, ValueError):
        overlay_opacity = 0.46
    overlay_opacity = max(0.0, min(0.82, overlay_opacity))

    defaults = {
        'title_color': title_color,
        'content_color': content_color,
        'accent_color': accent_color,
        'anchor_pad_color': anchor_pad_color,
    }

    return {
        'title': title,
        'role': slide_data.get('role') or ('cover' if index == 0 else 'content'),
        'render_mode': str(slide_data.get('render_mode') or '').strip().lower(),
        'visual_metaphor': slide_data.get('visual_metaphor') or '',
        'meaning_map': slide_data.get('meaning_map') if isinstance(slide_data.get('meaning_map'), dict) else {},
        'layout': layout,
        'body_blocks': body_blocks,
        'text_anchors': _coerce_text_anchors(slide_data, defaults),
        'bg_color': bg_color,
        'title_color': title_color,
        'content_color': content_color,
        'accent_color': accent_color,
        'anchor_pad_color': anchor_pad_color,
        'overlay_color': overlay_color,
        'overlay_opacity': overlay_opacity,
        'background_image': slide_data.get('background_image'),
        'footer': slide_data.get('footer') or deck_defaults.get('footer') or '',
    }


def _add_overlay(slide, box, color_hex, opacity):
    x, y, w, h = box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    _set_shape_opacity(shape, opacity)
    return shape


def _set_shape_opacity(shape, opacity):
    opacity = _clamp_float(opacity, 0.70, 0.0, 1.0)
    transparency = int(round((1.0 - opacity) * 100))
    try:
        shape.fill.transparency = transparency
    except Exception:
        shape.fill.transparency = float(transparency)


def _normalized_box_to_inches(box):
    x, y, w, h = box
    return (
        x * 13.333,
        y * 7.5,
        w * 13.333,
        h * 7.5,
    )


def _normalized_point_to_inches(point):
    if not isinstance(point, dict):
        return None
    try:
        x = _clamp_float(point.get('x'), 0.5, 0.0, 1.0) * 13.333
        y = _clamp_float(point.get('y'), 0.5, 0.0, 1.0) * 7.5
    except Exception:
        return None
    return x, y


def _has_structural_anchors(slide_info):
    anchors = slide_info.get('text_anchors') or []
    content_valid = 0
    evidence_valid = 0
    for anchor in anchors:
        if not anchor.get('text') or not anchor.get('bind_to') or anchor.get('type') == 'title':
            continue
        content_valid += 1
        if anchor.get('evidence_ref'):
            evidence_valid += 1
    return content_valid >= 2 and (evidence_valid >= 2 or content_valid >= 4)


def _synthesize_anchor_connector(anchor):
    if not anchor.get('connector_required'):
        return None
    x, y, w, h = anchor['box']
    target = anchor.get('target_point')
    if isinstance(target, dict):
        start_x = _clamp_float(target.get('x'), x + w / 2, 0.0, 1.0)
        start_y = _clamp_float(target.get('y'), y + h / 2, 0.0, 1.0)
    else:
        start_x = x - 0.04 if x > 0.52 else x + w + 0.04
        start_y = y + h / 2
        start_x = max(0.03, min(0.97, start_x))
        start_y = max(0.05, min(0.95, start_y))

    end_x = x if start_x < x else x + w
    end_y = y + h / 2
    return {
        'from': {'x': start_x, 'y': start_y},
        'to': {'x': max(0.02, min(0.98, end_x)), 'y': max(0.02, min(0.98, end_y))},
    }


def _render_connector(slide, connector, color_hex):
    if not isinstance(connector, dict):
        return
    start = _normalized_point_to_inches(connector.get('from') or connector.get('start'))
    end = _normalized_point_to_inches(connector.get('to') or connector.get('end'))
    if not start or not end:
        return

    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start[0]),
        Inches(start[1]),
        Inches(end[0]),
        Inches(end[1]),
    )
    line.line.color.rgb = _rgb(color_hex)
    line.line.width = Pt(1.3)

    dot_size = 0.08
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(start[0] - dot_size / 2),
        Inches(start[1] - dot_size / 2),
        Inches(dot_size),
        Inches(dot_size),
    )
    dot.line.color.rgb = _rgb(color_hex)
    dot.fill.solid()
    dot.fill.fore_color.rgb = _rgb(color_hex)


def _render_anchor_text(slide, anchor, slide_info):
    x, y, w, h = _normalized_box_to_inches(anchor['box'])
    if anchor['pad']:
        pad = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        pad.line.color.rgb = _rgb(anchor['pad_color'])
        pad.line.transparency = 70
        pad.fill.solid()
        pad.fill.fore_color.rgb = _rgb(anchor['pad_color'])
        _set_shape_opacity(pad, anchor['pad_opacity'])

    textbox = _add_textbox(slide, (x + 0.03, y + 0.02, max(0.1, w - 0.06), max(0.1, h - 0.04)))
    tf = textbox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    anchor_type = anchor['type']
    is_title = anchor_type == 'title'
    is_stat = anchor_type in ('stat', 'metric', 'number')
    dense_anchor_count = sum(
        1 for item in slide_info.get('text_anchors', [])
        if item.get('type') != 'title' and item.get('text') and item.get('bind_to')
    )
    dense_mode = dense_anchor_count >= 4
    text_size = anchor['font_size']
    if is_stat:
        text_size = max(text_size, 20 if dense_mode else 24)
    if dense_mode and not is_title:
        text_size = min(text_size, 20 if is_stat else 16)

    _paragraph(
        tf,
        anchor['text'],
        text_size,
        anchor['color'],
        bold=is_title or is_stat,
        align=anchor['align'],
        space_after=3 if anchor.get('supporting_text') else 0,
    )
    if anchor.get('supporting_text'):
        supporting_size = max(9, min(text_size - 4, text_size * 0.72))
        _paragraph(
            tf,
            anchor['supporting_text'],
            supporting_size,
            anchor['supporting_color'],
            bold=False,
            align=anchor['align'],
            space_after=0,
        )


def _render_anchor_slide(slide, slide_info, index, total):
    anchors = list(slide_info['text_anchors'])
    has_title_anchor = any(anchor['type'] == 'title' for anchor in anchors)
    if not has_title_anchor and slide_info['title']:
        anchors.insert(0, {
            'id': 'title_auto',
            'type': 'title',
            'text': slide_info['title'],
            'supporting_text': '',
            'box': (0.06, 0.06, 0.48, 0.12),
            'align': PP_ALIGN.LEFT,
            'font_size': 30,
            'color': slide_info['title_color'],
            'supporting_color': slide_info['content_color'],
            'pad': False,
            'pad_color': slide_info['anchor_pad_color'],
            'pad_opacity': 0.0,
            'connector': None,
        })

    for anchor in anchors:
        connector = anchor.get('connector') or _synthesize_anchor_connector(anchor)
        _render_connector(slide, connector, slide_info['accent_color'])
    for anchor in anchors:
        _render_anchor_text(slide, anchor, slide_info)
    _render_footer(slide, slide_info, index, total)


def _estimate_body_size(blocks):
    total = sum(len(str(block.get('text') or '')) + sum(len(str(i)) for i in block.get('items', []))
                for block in blocks if isinstance(block, dict))
    if total > 520:
        return 15
    if total > 340:
        return 17
    return 19


def _render_body_blocks(slide, box, slide_info, align):
    body_box = (box[0], box[1] + 1.28, box[2], max(0.7, box[3] - 1.42))
    body_shape = _add_textbox(slide, body_box)
    tf = body_shape.text_frame
    body_size = _estimate_body_size(slide_info['body_blocks'])

    for block in slide_info['body_blocks'][:7]:
        if not isinstance(block, dict):
            block = {'type': 'bullet', 'text': block}
        block_type = str(block.get('type') or 'bullet').lower()
        text = str(block.get('text') or '').strip()
        items = block.get('items') if isinstance(block.get('items'), list) else []

        if block_type in ('kicker', 'eyebrow'):
            _paragraph(tf, text, 13, slide_info['accent_color'], bold=True, align=align, space_after=6)
        elif block_type in ('stat', 'metric'):
            value = str(block.get('value') or text).strip()
            label = str(block.get('label') or '').strip()
            _paragraph(tf, value, min(30, body_size + 10), slide_info['accent_color'], bold=True, align=align, space_after=2)
            if label:
                _paragraph(tf, label, max(12, body_size - 2), slide_info['content_color'], align=align, space_after=10)
        elif items:
            if text:
                _paragraph(tf, text, body_size, slide_info['content_color'], bold=True, align=align, space_after=5)
            for item in items[:5]:
                _paragraph(tf, f'- {item}', body_size - 1, slide_info['content_color'], align=align, space_after=4)
        else:
            prefix = '- ' if block_type in ('bullet', 'point') else ''
            _paragraph(tf, f'{prefix}{text}', body_size, slide_info['content_color'], align=align, space_after=7)


def _render_footer(slide, slide_info, index, total):
    marker = f'{index + 1:02d}/{total:02d}'
    footer_text = slide_info['footer']
    if footer_text:
        marker = f'{footer_text}  |  {marker}'
    footer_shape = _add_textbox(slide, (0.78, 6.92, 11.80, 0.28), vertical_anchor=MSO_ANCHOR.MIDDLE)
    _paragraph(footer_shape.text_frame, marker, 9, slide_info['content_color'], align=PP_ALIGN.RIGHT, space_after=0)


def _render_slide(slide, prs, slide_info, index, total, hybrid_mode):
    _set_slide_background(slide, slide_info['bg_color'])
    has_bg_image = bool(hybrid_mode and slide_info.get('background_image') and _add_background_image(prs, slide, slide_info['background_image']))
    use_anchor_mode = bool(
        has_bg_image
        and slide_info.get('text_anchors')
        and _has_structural_anchors(slide_info)
        and slide_info.get('render_mode') != 'fallback_layout'
    )

    if use_anchor_mode:
        _render_anchor_slide(slide, slide_info, index, total)
        return

    box = LAYOUT_BOXES[slide_info['layout']]
    align = PP_ALIGN.CENTER if slide_info['layout'] == 'center' else PP_ALIGN.LEFT
    if has_bg_image or slide_info['layout'] in ('center', 'full'):
        _add_overlay(slide, box, slide_info['overlay_color'], slide_info['overlay_opacity'])

    title_shape = _add_textbox(slide, (box[0], box[1], box[2], 1.10), vertical_anchor=MSO_ANCHOR.MIDDLE)
    title_len = len(slide_info['title'])
    title_size = 42 if index == 0 and title_len < 24 else 34
    if title_len > 42:
        title_size = 28
    _paragraph(title_shape.text_frame, slide_info['title'], title_size, slide_info['title_color'], bold=True, align=align, space_after=4)

    _render_body_blocks(slide, box, slide_info, align)

    _render_footer(slide, slide_info, index, total)


def create_pptx_from_json(json_data, output_path, hybrid_mode=False):
    """
    Render an editable PPTX using blank slides, cover backgrounds, and CJK-safe text boxes.
    Supports both the new design schema and the legacy title/content/colors schema.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    slides_data = json_data.get('slides') or []
    if not slides_data:
        slides_data = [{
            'title': '簡報生成結果',
            'content': ['目前沒有可用頁面資料，請重新提交提示詞或素材。'],
            'layout': 'center',
            'bg_color': '#111827',
            'title_color': '#FFFFFF',
            'content_color': '#E5E7EB',
        }]

    design_system = json_data.get('design_system') if isinstance(json_data.get('design_system'), dict) else {}
    deck_brief = json_data.get('deck_brief') if isinstance(json_data.get('deck_brief'), dict) else {}
    deck_defaults = {
        'colors': design_system.get('colors') if isinstance(design_system.get('colors'), dict) else {},
        'footer': deck_brief.get('audience') or deck_brief.get('purpose') or ''
    }

    total = len(slides_data)
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        slide_info = _normalize_slide(slide_data, i, deck_defaults)
        _render_slide(slide, prs, slide_info, i, total, hybrid_mode)

    prs.save(output_path)


def _strip_data_uri_prefix(img_str):
    """剝離 data:image/xxx;base64, 前綴，只留 base64 本體"""
    if ',' in img_str:
        return img_str.split(',', 1)[1]
    return img_str


def _collect_uploaded_images(files):
    all_base64_images = []
    processed_any_file = False

    for file in files:
        if not file or not file.filename:
            continue
        if not allowed_file(file.filename):
            continue
        processed_any_file = True
        safe_filename = _safe_upload_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)

        try:
            file.save(filepath)
            images = process_file_to_images(filepath)
            for img in images:
                all_base64_images.append(_strip_data_uri_prefix(img))
        except Exception as e:
            original_name = secure_filename(file.filename) or file.filename
            raise ValueError(f"素材解析失敗：{original_name}") from e
        finally:
            try:
                os.remove(filepath)
            except Exception:
                pass

    return all_base64_images, processed_any_file


def _format_vault_error(resp, fallback):
    try:
        data = resp.json()
        detail = data.get('detail') or data.get('error') or data.get('message')
        if detail:
            return str(detail)
    except Exception:
        pass
    text = (resp.text or '').replace('\n', ' ').strip()
    return text[:240] if text else fallback


def _compact_text(text):
    return re.sub(r'\s+', '', str(text or '').replace('–', '-').replace('—', '-'))


def _money_to_float(value):
    text = re.sub(r'[^\d.]', '', str(value or ''))
    return float(text) if text else 0.0


def _money_to_int(value):
    return int(round(_money_to_float(value)))


def _format_money(value):
    return f"{_money_to_int(value):,}"


def _pdf_page_texts(doc):
    return [doc.load_page(i).get_text('text') for i in range(len(doc))]


def _find_page_by_terms(page_texts, terms):
    compact_terms = [_compact_text(term) for term in terms]
    for index, text in enumerate(page_texts):
        compact = _compact_text(text)
        if all(term in compact for term in compact_terms):
            return index
    return None


def _numeric_lines_after_label(text, label, limit=12):
    lines = [line.strip() for line in str(text or '').splitlines() if line.strip()]
    for index, line in enumerate(lines):
        if line == str(label):
            values = []
            for candidate in lines[index + 1:]:
                if len(values) >= limit:
                    break
                if re.fullmatch(r'\d{1,3}(?:,\d{3})*(?:\.\d+)?|\d+(?:\.\d+)?', candidate):
                    values.append(candidate)
                elif values:
                    break
            if values:
                return values
    return []


def _median_number(values, fallback=0.0):
    cleaned = sorted(float(value) for value in values if value is not None)
    if not cleaned:
        return fallback
    middle = len(cleaned) // 2
    if len(cleaned) % 2:
        return cleaned[middle]
    return (cleaned[middle - 1] + cleaned[middle]) / 2


def _ocr_numeric_value(text):
    normalized = str(text or '').strip().translate(str.maketrans({
        'O': '0',
        'o': '0',
        'I': '1',
        'l': '1',
        '|': '1',
    }))
    normalized = re.sub(r'[^\d,]', '', normalized)
    if not normalized or not re.fullmatch(r'\d{1,3}(?:,\d{3})*|\d+', normalized):
        return None
    return int(normalized.replace(',', ''))


def _parse_ocr_confidence(value):
    try:
        return float(value)
    except (TypeError, ValueError):
        return -1.0


def _merge_rects(rects):
    rects = [rect for rect in rects if rect is not None]
    if not rects:
        return None
    merged = fitz.Rect(rects[0])
    for rect in rects[1:]:
        merged.include_rect(rect)
    return merged


def _ocr_page_numeric_words(page):
    if pytesseract is None:
        raise ValueError('OCR 套件未安裝，無法安全定位 GF 提款表格。')

    pix = page.get_pixmap(matrix=fitz.Matrix(GF_OCR_ZOOM, GF_OCR_ZOOM), alpha=False)
    image = Image.open(io.BytesIO(pix.tobytes('png')))
    data = pytesseract.image_to_data(
        image,
        lang=GF_OCR_LANG,
        config='--psm 6',
        output_type=pytesseract.Output.DICT,
    )
    scale_x = page.rect.width / max(1, image.width)
    scale_y = page.rect.height / max(1, image.height)
    words = []
    for index, raw_text in enumerate(data.get('text', [])):
        text = str(raw_text or '').strip()
        value = _ocr_numeric_value(text)
        if value is None:
            continue
        confidence = _parse_ocr_confidence(data.get('conf', [])[index])
        if 0 <= confidence < GF_OCR_MIN_CONFIDENCE:
            continue
        x0 = page.rect.x0 + int(data['left'][index]) * scale_x
        y0 = page.rect.y0 + int(data['top'][index]) * scale_y
        x1 = x0 + int(data['width'][index]) * scale_x
        y1 = y0 + int(data['height'][index]) * scale_y
        words.append({
            'text': text,
            'value': value,
            'confidence': confidence,
            'rect': fitz.Rect(x0, y0, x1, y1),
        })
    return words


def _cluster_ocr_words_by_y(words, tolerance=4.0):
    clusters = []
    for word in sorted(words, key=lambda item: ((item['rect'].y0 + item['rect'].y1) / 2, item['rect'].x0)):
        center_y = (word['rect'].y0 + word['rect'].y1) / 2
        for cluster in clusters:
            if abs(cluster['center_y'] - center_y) <= tolerance:
                cluster['words'].append(word)
                cluster['center_y'] = _median_number(
                    [(w['rect'].y0 + w['rect'].y1) / 2 for w in cluster['words']],
                    center_y,
                )
                break
        else:
            clusters.append({'center_y': center_y, 'words': [word]})
    return clusters


def _cluster_ocr_words_by_x(words, tolerance=10.0):
    clusters = []
    for word in sorted(words, key=lambda item: (item['rect'].x0 + item['rect'].x1) / 2):
        center_x = (word['rect'].x0 + word['rect'].x1) / 2
        for cluster in clusters:
            if abs(cluster['center_x'] - center_x) <= tolerance:
                cluster['words'].append(word)
                cluster['center_x'] = _median_number(
                    [(w['rect'].x0 + w['rect'].x1) / 2 for w in cluster['words']],
                    center_x,
                )
                break
        else:
            clusters.append({'center_x': center_x, 'words': [word]})
    return clusters


def _longest_consecutive_year_run(year_rows):
    years = sorted(year_rows)
    longest = current = 0
    previous = None
    for year in years:
        if previous is None or year == previous + 1:
            current += 1
        else:
            longest = max(longest, current)
            current = 1
        previous = year
    return max(longest, current)


def _build_year_rows_from_ocr_words(words):
    by_year = {}
    for word in words:
        by_year.setdefault(word['value'], []).append(word)

    year_rows = {}
    for year, year_words in by_year.items():
        clusters = _cluster_ocr_words_by_y(year_words)
        if not clusters:
            continue
        best = max(
            clusters,
            key=lambda cluster: (
                len(cluster['words']),
                _median_number([w['confidence'] for w in cluster['words']], -1),
            ),
        )
        year_rows[year] = {
            'center_y': best['center_y'],
            'rect': _merge_rects([word['rect'] for word in best['words']]),
            'count': len(best['words']),
        }
    return year_rows


def _ocr_year_row_map(page, max_year=45):
    left_limit = page.rect.width * 0.18
    year_words = [
        word for word in _ocr_page_numeric_words(page)
        if 1 <= word['value'] <= max_year and word['rect'].x0 <= left_limit
    ]
    column_candidates = []
    for cluster in _cluster_ocr_words_by_x(year_words):
        rows = _build_year_rows_from_ocr_words(cluster['words'])
        if not rows:
            continue
        step = _ocr_row_step(rows) if len(rows) > 1 else 0
        column_candidates.append({
            'center_x': cluster['center_x'],
            'rows': rows,
            'run': _longest_consecutive_year_run(rows),
            'count': len(rows),
            'step': step,
        })
    if not column_candidates:
        raise ValueError('OCR 未能定位 GF 提款表格年份欄。')

    best_column = max(
        column_candidates,
        key=lambda item: (
            item['run'],
            item['count'],
            1 if 8.0 <= item['step'] <= 16.0 else 0,
            item['center_x'],
        ),
    )
    year_rows = best_column['rows']
    if len(year_rows) < 10:
        raise ValueError('OCR 未能可靠定位 GF 提款表格年份列。')
    return year_rows


def _ocr_row_step(year_rows):
    ordered = sorted((year, row['center_y']) for year, row in year_rows.items())
    diffs = [
        y2 - y1
        for (_, y1), (_, y2) in zip(ordered, ordered[1:])
        if 5.0 <= y2 - y1 <= 30.0
    ]
    return _median_number(diffs, fallback=12.0)


def _ocr_row_bounds(year_rows, year):
    if year not in year_rows:
        raise ValueError(f'OCR 找不到第 {year} 年在 GF 提款表格的位置。')
    center_y = year_rows[year]['center_y']
    step = _ocr_row_step(year_rows)
    previous_rows = [row['center_y'] for row_year, row in year_rows.items() if row_year < year]
    next_rows = [row['center_y'] for row_year, row in year_rows.items() if row_year > year]
    top = (max(previous_rows) + center_y) / 2 if previous_rows else center_y - step / 2
    bottom = (min(next_rows) + center_y) / 2 if next_rows else center_y + step / 2
    padding = max(1.0, min(3.0, step * 0.16))
    return top - padding, bottom + padding


def _is_pdf_numeric_word(text):
    normalized = re.sub(r'[^\d,.]', '', str(text or ''))
    return bool(re.fullmatch(r'\d{1,3}(?:,\d{3})*(?:\.\d+)?|\d+(?:\.\d+)?', normalized))


def _pdf_numeric_cells_for_ocr_row(page, year_rows, year):
    top, bottom = _ocr_row_bounds(year_rows, year)
    cells = []
    for word in page.get_text('words'):
        x0, y0, x1, y1, text = word[:5]
        center_y = (y0 + y1) / 2
        if center_y < top or center_y > bottom or not _is_pdf_numeric_word(text):
            continue
        cells.append({
            'text': str(text),
            'value': _money_to_int(text),
            'rect': fitz.Rect(x0, y0, x1, y1),
        })
    cells.sort(key=lambda cell: (cell['rect'].x0, cell['rect'].y0))
    left_limit = page.rect.width * 0.18
    if not any(cell['value'] == year and cell['rect'].x0 <= left_limit for cell in cells):
        raise ValueError(f'PDF 抽取結果未能對上 OCR 定位的第 {year} 年列。')
    return cells


def _pdf_value_cells(cells, page):
    left_limit = page.rect.width * 0.18
    return [cell for cell in cells if cell['rect'].x0 > left_limit]


def _rightmost_pdf_value_cell(cells, page, label):
    value_cells = _pdf_value_cells(cells, page)
    if not value_cells:
        raise ValueError(f'找不到 {label} 的 PDF 金額欄位。')
    return value_cells[-1]


def _find_pdf_amount_cell(cells, page, amount, label):
    matches = [
        cell for cell in _pdf_value_cells(cells, page)
        if cell['value'] == amount
    ]
    if not matches:
        raise ValueError(f'PDF 第 {label} 列找不到每年提款金額 {_format_money(amount)}。')
    return matches[0]


def _extract_gf_withdrawal_from_ocr_pdf(doc, pages):
    target_year = GF_WITHDRAWAL_TARGET_YEAR
    amount_page = doc.load_page(pages['withdrawal_amount_table'])
    surrender_page = doc.load_page(pages['withdrawal_surrender_table'])
    amount_year_rows = _ocr_year_row_map(amount_page)
    surrender_year_rows = _ocr_year_row_map(surrender_page)

    yearly_withdrawals = {}
    for year in range(2, target_year + 1):
        row_cells = _pdf_numeric_cells_for_ocr_row(amount_page, amount_year_rows, year)
        withdrawal_cell = _rightmost_pdf_value_cell(row_cells, amount_page, f'第 {year} 年提款金額')
        yearly_withdrawals[year] = withdrawal_cell['value']

    positive_years = [year for year, amount in yearly_withdrawals.items() if amount > 0]
    if not positive_years:
        raise ValueError('PDF/OCR 驗證後找不到提款開始年份。')

    start_year = positive_years[0]
    annual_withdrawal = yearly_withdrawals[start_year]
    for year in range(start_year, target_year + 1):
        amount = yearly_withdrawals.get(year)
        if amount != annual_withdrawal:
            raise ValueError(
                f'第 {start_year} 至第 {target_year} 年提款金額不一致，'
                f'第 {year} 年是 {_format_money(amount or 0)}。'
            )

    total_withdrawal = sum(yearly_withdrawals[year] for year in range(start_year, target_year + 1))

    start_surrender_cells = _pdf_numeric_cells_for_ocr_row(surrender_page, surrender_year_rows, start_year)
    target_surrender_cells = _pdf_numeric_cells_for_ocr_row(surrender_page, surrender_year_rows, target_year)
    start_amount_cell = _find_pdf_amount_cell(start_surrender_cells, surrender_page, annual_withdrawal, start_year)
    target_amount_cell = _find_pdf_amount_cell(target_surrender_cells, surrender_page, annual_withdrawal, target_year)
    balance_cell = _rightmost_pdf_value_cell(target_surrender_cells, surrender_page, f'第 {target_year} 年戶口餘額')
    account_balance = balance_cell['value']
    if account_balance <= annual_withdrawal:
        raise ValueError(
            f'第 {target_year} 年戶口餘額 {_format_money(account_balance)} 異常，'
            '已停止生成 GF 簡報。'
        )

    surrender_crop = GF_CROP_RULES['withdrawal_surrender_table']['crop']
    start_box = _rect_to_crop_box(
        surrender_page,
        start_amount_cell['rect'],
        surrender_crop,
        pad_x=0.018,
        pad_y=0.006,
    )
    target_box = _rect_to_crop_box(
        surrender_page,
        target_amount_cell['rect'],
        surrender_crop,
        pad_x=0.018,
        pad_y=0.006,
    )
    balance_box = _rect_to_crop_box(
        surrender_page,
        balance_cell['rect'],
        surrender_crop,
        pad_x=0.012,
        pad_y=0.006,
    )
    highlight_regions = {
        'withdrawal_start_amount': start_box,
        'withdrawal_target_amount': target_box,
        'withdrawal_amount_range': _merge_crop_boxes(start_box, target_box, pad_x=0.006, pad_y=0.004),
        'account_balance_target': balance_box,
    }

    return {
        'withdrawal_start_year': str(start_year),
        'annual_withdrawal': _format_money(annual_withdrawal),
        'total_withdrawal': _format_money(total_withdrawal),
        'account_balance_at_target_year': _format_money(account_balance),
        '_withdrawal_highlight_regions': {
            key: value for key, value in highlight_regions.items() if value
        },
    }


def _extract_summary_values(summary_text):
    row_15 = _numeric_lines_after_label(summary_text, 15, limit=10)
    row_25 = _numeric_lines_after_label(summary_text, 25, limit=10)
    if len(row_15) < 6 or len(row_25) < 5:
        raise ValueError('找不到基本計劃摘要中的 15 年 / 25 年數值。')

    paid_premium_total = row_15[5]
    summary_15y = row_15[4]
    summary_25y = row_25[4]
    paid = max(1.0, _money_to_float(paid_premium_total))
    growth_15y = int((_money_to_float(summary_15y) - paid) / paid * 100)
    growth_25y = int((_money_to_float(summary_25y) - paid) / paid * 100)

    return {
        'paid_premium_total': _format_money(paid_premium_total),
        'summary_15y': _format_money(summary_15y),
        'summary_25y': _format_money(summary_25y),
        'growth_15y': str(growth_15y),
        'growth_25y': str(growth_25y),
    }


def _extract_annual_premium(first_page_text):
    match = re.search(r'環宇盈活儲蓄保險計劃（5 年\s*繳費）\s*([\d,]+)', first_page_text, re.S)
    if match:
        return _format_money(match.group(1))
    match = re.search(r'投保時年繳總保費：\s*([\d,]+(?:\.\d+)?)', first_page_text)
    return _format_money(match.group(1)) if match else ''


def _extract_currency(first_page_text):
    match = re.search(r'保單貨幣：\s*([^\n\r]+)', first_page_text)
    value = match.group(1).strip() if match else ''
    if '美元' in value:
        return 'USD'
    if '港' in value:
        return 'HKD'
    return value or 'USD'


def _extract_withdrawal_year_amounts(withdrawal_text):
    year_amounts = []
    for match in re.finditer(r'\n([\d,]+)\n(\d{1,2})\n', withdrawal_text):
        amount = _money_to_int(match.group(1))
        year = int(match.group(2))
        if amount > 0 and 2 <= year <= GF_WITHDRAWAL_TARGET_YEAR:
            year_amounts.append((year, amount))
    return year_amounts


def _extract_withdrawal_start_and_amount(withdrawal_text):
    year_amounts = _extract_withdrawal_year_amounts(withdrawal_text)
    if year_amounts:
        year, amount = year_amounts[0]
        return str(year), _format_money(amount)
    return '', ''


def _extract_account_balance_at_year(surrender_text, target_year):
    values = _numeric_lines_after_label(surrender_text, target_year, limit=3)
    if not values:
        raise ValueError(f'找不到第 {target_year} 年的提款後戶口餘額。')
    return _format_money(values[0])


def _search_value_rect(page, value, occurrence=0):
    rects = page.search_for(str(value or ''))
    if not rects or occurrence >= len(rects):
        return None
    return rects[occurrence]


def _rect_to_crop_box(page, rect, crop, pad_x=0.012, pad_y=0.006):
    if rect is None:
        return None
    page_rect = page.rect
    crop_x0 = page_rect.x0 + page_rect.width * crop[0]
    crop_y0 = page_rect.y0 + page_rect.height * crop[1]
    crop_w = page_rect.width * (crop[2] - crop[0])
    crop_h = page_rect.height * (crop[3] - crop[1])
    x = (rect.x0 - crop_x0) / crop_w
    y = (rect.y0 - crop_y0) / crop_h
    w = rect.width / crop_w
    h = rect.height / crop_h
    x = max(0.0, x - pad_x)
    y = max(0.0, y - pad_y)
    w = min(1.0 - x, w + pad_x * 2)
    h = min(1.0 - y, h + pad_y * 2)
    return (x, y, w, h)


def _merge_crop_boxes(box_a, box_b, pad_x=0.01, pad_y=0.006):
    if not box_a or not box_b:
        return box_a or box_b
    x0 = max(0.0, min(box_a[0], box_b[0]) - pad_x)
    y0 = max(0.0, min(box_a[1], box_b[1]) - pad_y)
    x1 = min(1.0, max(box_a[0] + box_a[2], box_b[0] + box_b[2]) + pad_x)
    y1 = min(1.0, max(box_a[1] + box_a[3], box_b[1] + box_b[3]) + pad_y)
    return (x0, y0, x1 - x0, y1 - y0)


def _extract_highlight_regions(doc, gf_data):
    regions = {}
    summary_page = doc.load_page(gf_data['pages']['summary_table'])
    summary_crop = GF_CROP_RULES['summary_table']['crop']
    regions['summary_15y'] = _rect_to_crop_box(
        summary_page,
        _search_value_rect(summary_page, gf_data['summary_15y'], 0),
        summary_crop,
    )
    regions['summary_25y'] = _rect_to_crop_box(
        summary_page,
        _search_value_rect(summary_page, gf_data['summary_25y'], 0),
        summary_crop,
    )

    if gf_data.get('has_withdrawal'):
        verified_regions = gf_data.get('_withdrawal_highlight_regions') or {}
        if verified_regions:
            regions.update(verified_regions)
            return {key: value for key, value in regions.items() if value}

        surrender_page = doc.load_page(gf_data['pages']['withdrawal_surrender_table'])
        surrender_crop = GF_CROP_RULES['withdrawal_surrender_table']['crop']
        start_offset = 0
        target_offset = max(
            0,
            int(gf_data['withdrawal_target_year']) - int(gf_data['withdrawal_start_year'])
        )
        start_box = _rect_to_crop_box(
            surrender_page,
            _search_value_rect(surrender_page, gf_data['annual_withdrawal'], start_offset),
            surrender_crop,
            pad_x=0.018,
            pad_y=0.006,
        )
        target_withdrawal_box = _rect_to_crop_box(
            surrender_page,
            _search_value_rect(surrender_page, gf_data['annual_withdrawal'], target_offset),
            surrender_crop,
            pad_x=0.018,
            pad_y=0.006,
        )
        regions['withdrawal_amount_range'] = _merge_crop_boxes(
            start_box,
            target_withdrawal_box,
            pad_x=0.006,
            pad_y=0.004,
        )
        regions['withdrawal_start_amount'] = start_box
        regions['withdrawal_target_amount'] = target_withdrawal_box
        regions['account_balance_target'] = _rect_to_crop_box(
            surrender_page,
            _search_value_rect(surrender_page, gf_data['account_balance_at_target_year'], 0),
            surrender_crop,
            pad_x=0.012,
            pad_y=0.006,
        )
    return {key: value for key, value in regions.items() if value}


def _parse_gf_proposal(pdf_path, client_name, agent_name, withdrawal_mode):
    if not os.path.exists(pdf_path):
        raise ValueError('找不到上傳的建議書 PDF。')

    doc = fitz.open(pdf_path)
    try:
        page_texts = _pdf_page_texts(doc)
        first_page_text = page_texts[0] if page_texts else ''
        full_text = '\n'.join(page_texts)
        if '環宇盈活儲蓄保險計劃' not in full_text:
            raise ValueError('這份文件不像 GF 建議書：找不到「環宇盈活儲蓄保險計劃」。')

        pages = {}
        for key, rule in GF_CROP_RULES.items():
            pages[key] = _find_page_by_terms(page_texts, rule['terms'])

        if pages['proposal_summary'] is None:
            raise ValueError('找不到「建議書摘要 / 保障摘要」頁面。')
        if pages['summary_table'] is None:
            raise ValueError('找不到「基本計劃 – 說明摘要」頁面。')

        has_withdrawal = pages['withdrawal_amount_table'] is not None and pages['withdrawal_surrender_table'] is not None
        if withdrawal_mode == 'yes' and not has_withdrawal:
            raise ValueError('建議書未包含可用的提款例子。')
        if withdrawal_mode == 'no':
            has_withdrawal = False

        summary = _extract_summary_values(page_texts[pages['summary_table']])
        data = {
            'client_name': client_name,
            'agent_name': agent_name,
            'plan_name': '環宇盈活儲蓄保險計劃（5 年繳費）',
            'currency': _extract_currency(first_page_text),
            'annual_premium': _extract_annual_premium(first_page_text),
            'withdrawal_target_year': str(GF_WITHDRAWAL_TARGET_YEAR),
            'has_withdrawal': has_withdrawal,
            'pages': pages,
            **summary,
        }

        if has_withdrawal:
            data.update(_extract_gf_withdrawal_from_ocr_pdf(doc, pages))
        else:
            data.update({
                'withdrawal_start_year': '',
                'annual_withdrawal': '',
                'total_withdrawal': '',
                'account_balance_at_target_year': '',
            })

        data['highlight_regions'] = _extract_highlight_regions(doc, data)
        return data
    finally:
        doc.close()


def _render_pdf_crop(pdf_path, page_index, crop, output_path):
    doc = fitz.open(pdf_path)
    try:
        page = doc.load_page(page_index)
        rect = page.rect
        x0, y0, x1, y1 = crop
        clip = fitz.Rect(
            rect.x0 + rect.width * x0,
            rect.y0 + rect.height * y0,
            rect.x0 + rect.width * x1,
            rect.y0 + rect.height * y1,
        )
        pix = page.get_pixmap(matrix=fitz.Matrix(2.4, 2.4), clip=clip, alpha=False)
        pix.save(output_path)
        return output_path
    finally:
        doc.close()


def _render_gf_crops(pdf_path, gf_data, task_id):
    crop_paths = {}
    for key in ('proposal_summary', 'summary_table', 'withdrawal_surrender_table'):
        if key == 'withdrawal_surrender_table' and not gf_data.get('has_withdrawal'):
            continue
        page_index = gf_data['pages'].get(key)
        if page_index is None:
            continue
        crop_path = os.path.join(app.config['UPLOAD_FOLDER'], f'gf_{task_id}_{key}.png')
        _render_pdf_crop(pdf_path, page_index, GF_CROP_RULES[key]['crop'], crop_path)
        crop_paths[key] = crop_path
    return crop_paths


def _slide_at(prs, one_based_index):
    if len(prs.slides) < one_based_index:
        raise ValueError(f'GF 模板缺少第 {one_based_index} 頁。')
    return prs.slides[one_based_index - 1]


def _delete_slide(prs, zero_based_index):
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[zero_based_index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def _replace_text_in_slide(slide, replacements):
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, str(new))


def _set_text_shape_text(shape, text):
    if not getattr(shape, 'has_text_frame', False):
        return
    tf = shape.text_frame
    first_paragraph = tf.paragraphs[0] if tf.paragraphs else None
    first_run = first_paragraph.runs[0] if first_paragraph and first_paragraph.runs else None
    align = first_paragraph.alignment if first_paragraph else None
    font_name = first_run.font.name if first_run else None
    font_size = first_run.font.size if first_run else None
    font_bold = first_run.font.bold if first_run else None
    font_color = None
    try:
        font_color = first_run.font.color.rgb if first_run else None
    except Exception:
        font_color = None

    tf.clear()
    lines = str(text).splitlines() or ['']
    for index, line_text in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        if align is not None:
            paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line_text
        if font_name:
            run.font.name = font_name
            r_pr = run._r.get_or_add_rPr()
            for tag in ('a:latin', 'a:ea', 'a:cs'):
                font_elem = r_pr.find(qn(tag))
                if font_elem is None:
                    font_elem = OxmlElement(tag)
                    r_pr.append(font_elem)
                font_elem.set('typeface', font_name)
        if font_size:
            run.font.size = font_size
        if font_bold is not None:
            run.font.bold = font_bold
        if font_color:
            run.font.color.rgb = font_color


def _find_text_shape(slide, contains):
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and contains in shape.text:
            return shape
    return None


def _find_text_shape_with_all(slide, terms):
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        text = shape.text
        if all(term in text for term in terms):
            return shape
    return None


def _remove_text_shapes(slide, exact_texts):
    targets = {' '.join(text.split()) for text in exact_texts}
    to_remove = []
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        normalized = ' '.join(shape.text.split())
        if normalized in targets:
            to_remove.append(shape)
    for shape in to_remove:
        shape._element.getparent().remove(shape._element)


def _shape_color(shape, attr):
    try:
        color = getattr(shape, attr).color.rgb if attr == 'line' else shape.fill.fore_color.rgb
        return str(color) if color else None
    except Exception:
        return None


def _remove_gf_highlight_shapes(slide):
    to_remove = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and shape.text.strip():
            continue
        line_color = _shape_color(shape, 'line')
        fill_color = _shape_color(shape, 'fill')
        if line_color in {'FF0000', 'F5913F'} or fill_color == 'FF0000':
            to_remove.append(shape)
    for shape in to_remove:
        shape._element.getparent().remove(shape._element)


def _crop_box_to_slide_box(picture_shape, crop_box):
    x, y, w, h = crop_box
    return (
        int(picture_shape.left + picture_shape.width * x),
        int(picture_shape.top + picture_shape.height * y),
        int(picture_shape.width * w),
        int(picture_shape.height * h),
    )


def _add_highlight_box(slide, box, color_hex='FF0000', width_pt=1.5):
    left, top, width, height = box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    shape.line.color.rgb = _rgb(color_hex)
    shape.line.width = Pt(width_pt)
    return shape


def _add_emu_connector(slide, start_x, start_y, end_x, end_y, color_hex='F5913F', width_pt=1.4):
    if start_x == end_x and start_y == end_y:
        return None
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    line.line.color.rgb = _rgb(color_hex)
    line.line.width = Pt(width_pt)
    return line


def _add_emu_line(slide, start_x, start_y, end_x, end_y, color_hex='F5913F', width_pt=1.4):
    if end_x <= start_x:
        return None
    return _add_emu_connector(slide, start_x, start_y, end_x, end_y, color_hex, width_pt)


def _add_emu_elbow_line(slide, start_x, start_y, end_x, end_y, route_y, color_hex='F5913F', width_pt=1.4):
    if end_x <= start_x:
        return None
    route_y = int(route_y)
    if abs(route_y - start_y) < 2 and abs(route_y - end_y) < 2:
        return _add_emu_line(slide, start_x, start_y, end_x, end_y, color_hex, width_pt)

    created = []
    for sx, sy, ex, ey in (
        (start_x, start_y, start_x, route_y),
        (start_x, route_y, end_x, route_y),
        (end_x, route_y, end_x, end_y),
    ):
        line = _add_emu_connector(slide, int(sx), int(sy), int(ex), int(ey), color_hex, width_pt)
        if line:
            created.append(line)
    return created[-1] if created else None


def _center_shape_on_y(shape, center_y):
    shape.top = int(center_y - shape.height / 2)


def _apply_summary_highlights(slide, picture_shape, gf_data):
    highlights = gf_data.get('highlight_regions') or {}
    rows = [
        ('summary_15y', _find_text_shape(slide, '15年')),
        ('summary_25y', _find_text_shape(slide, '25年')),
    ]
    _remove_gf_highlight_shapes(slide)
    for key, label_shape in rows:
        if key not in highlights:
            continue
        box = _crop_box_to_slide_box(picture_shape, highlights[key])
        highlight = _add_highlight_box(slide, box)
        center_y = highlight.top + highlight.height / 2
        if label_shape:
            _center_shape_on_y(label_shape, center_y)
            start_x = int(label_shape.left + label_shape.width + Inches(1.05))
            end_x = int(highlight.left - Inches(0.08))
            _add_emu_line(slide, start_x, int(center_y), end_x, int(center_y))


def _apply_withdrawal_highlights(slide, picture_shape, gf_data):
    highlights = gf_data.get('highlight_regions') or {}
    _remove_gf_highlight_shapes(slide)

    range_box = highlights.get('withdrawal_amount_range')
    start_box = highlights.get('withdrawal_start_amount')
    balance_box = highlights.get('account_balance_target')

    if range_box:
        amount_highlight = _add_highlight_box(slide, _crop_box_to_slide_box(picture_shape, range_box))
        start_label = _find_text_shape(slide, '從第')
        if start_label and start_box:
            start_value_box = _crop_box_to_slide_box(picture_shape, start_box)
            start_center_y = start_value_box[1] + start_value_box[3] / 2
            _center_shape_on_y(start_label, start_center_y)
            title_shape = _find_text_shape_with_all(slide, ('每年提取', '戶口價值例子'))
            if title_shape:
                min_top = int(title_shape.top + title_shape.height + Inches(0.08))
                if start_label.top < min_top:
                    start_label.top = min_top
            _add_emu_line(
                slide,
                int(start_label.left + start_label.width + Inches(0.18)),
                int(start_label.top + start_label.height / 2),
                int(amount_highlight.left - Inches(0.08)),
                int(start_center_y),
            )

    if balance_box:
        balance_highlight = _add_highlight_box(slide, _crop_box_to_slide_box(picture_shape, balance_box))
        center_y = balance_highlight.top + balance_highlight.height / 2
        box_bottom_y = int(balance_highlight.top + balance_highlight.height)
        bottom_anchor_y = int(box_bottom_y + Inches(0.02))
        try:
            target_year = int(gf_data.get('withdrawal_target_year') or GF_WITHDRAWAL_TARGET_YEAR)
        except (TypeError, ValueError):
            target_year = GF_WITHDRAWAL_TARGET_YEAR
        try:
            start_year = int(gf_data.get('withdrawal_start_year') or target_year)
        except (TypeError, ValueError):
            start_year = target_year
        row_height = balance_highlight.height / 2
        if range_box:
            range_slide_box = _crop_box_to_slide_box(picture_shape, range_box)
            highlighted_rows = max(1, target_year - start_year + 1)
            row_height = max(1, range_slide_box[3] / highlighted_rows)
        rows_after_target = max(0, 30 - target_year)
        safe_line_y = int(box_bottom_y + row_height * (rows_after_target + 0.55))
        safe_line_y = max(safe_line_y, int(box_bottom_y + Inches(0.16)))
        picture_bottom_y = int(picture_shape.top + picture_shape.height - Inches(0.05))
        if safe_line_y > picture_bottom_y:
            safe_line_y = max(int(box_bottom_y + Inches(0.16)), picture_bottom_y)
        total_label = _find_text_shape(slide, '到第')
        right_label = _find_text_shape(slide, '戶口餘額$')
        if total_label:
            _center_shape_on_y(total_label, center_y)
            _add_emu_elbow_line(
                slide,
                int(total_label.left + total_label.width + Inches(0.15)),
                int(center_y),
                int(balance_highlight.left - Inches(0.08)),
                bottom_anchor_y,
                safe_line_y,
            )
        if right_label:
            _center_shape_on_y(right_label, center_y)
            _add_emu_elbow_line(
                slide,
                int(balance_highlight.left + balance_highlight.width + Inches(0.08)),
                bottom_anchor_y,
                int(right_label.left - Inches(0.12)),
                int(center_y),
                safe_line_y,
            )


def _find_largest_picture(slide):
    pictures = [shape for shape in slide.shapes if int(shape.shape_type) == 13]
    if not pictures:
        raise ValueError('GF 模板頁面找不到可替換的截圖位置。')
    return max(pictures, key=lambda shape: shape.width * shape.height)


def _replace_picture_keep_z(slide, picture_shape, image_path):
    sp_tree = slide.shapes._spTree
    old_element = picture_shape._element
    old_index = list(sp_tree).index(old_element)
    new_picture = slide.shapes.add_picture(
        image_path,
        picture_shape.left,
        picture_shape.top,
        width=picture_shape.width,
        height=picture_shape.height,
    )
    new_element = new_picture._element
    sp_tree.remove(new_element)
    sp_tree.insert(old_index, new_element)
    sp_tree.remove(old_element)
    return new_picture


def _finalize_gf_pptx(gf_data, crop_paths, output_path):
    if not os.path.exists(GF_TEMPLATE_PATH):
        raise ValueError('找不到 GF PPT 模板，請確認 assets/gf/gf_template.pptx 已部署。')

    prs = Presentation(GF_TEMPLATE_PATH)

    _replace_text_in_slide(_slide_at(prs, 1), {
        'Mr. Chung': gf_data['client_name'],
        'Prepared by Henry Chu': f"Prepared by {gf_data['agent_name']}",
    })

    slide_8 = _slide_at(prs, 8)
    _replace_picture_keep_z(slide_8, _find_largest_picture(slide_8), crop_paths['proposal_summary'])

    slide_9 = _slide_at(prs, 9)
    _remove_text_shapes(slide_9, {'一次過提取 戶口價值例子'})
    summary_picture = _replace_picture_keep_z(slide_9, _find_largest_picture(slide_9), crop_paths['summary_table'])
    label_15 = _find_text_shape(slide_9, '15年')
    label_25 = _find_text_shape(slide_9, '25年')
    if label_15:
        _set_text_shape_text(label_15, f"15年：{gf_data['growth_15y']}% 增長")
    if label_25:
        _set_text_shape_text(label_25, f"25年： {gf_data['growth_25y']}% 增長")
    _apply_summary_highlights(slide_9, summary_picture, gf_data)

    if gf_data.get('has_withdrawal'):
        slide_10 = _slide_at(prs, 10)
        withdrawal_picture = _replace_picture_keep_z(
            slide_10,
            _find_largest_picture(slide_10),
            crop_paths['withdrawal_surrender_table']
        )
        start_label = _find_text_shape(slide_10, '從第')
        total_label = _find_text_shape(slide_10, '到第')
        right_balance_label = _find_text_shape(slide_10, '戶口餘額$')
        if start_label:
            _set_text_shape_text(
                start_label,
                f"從第{gf_data['withdrawal_start_year']}年起\n每年提取${gf_data['annual_withdrawal']}"
            )
        if total_label:
            _set_text_shape_text(
                total_label,
                f"到第{gf_data['withdrawal_target_year']}年\n總提取${gf_data['total_withdrawal']}"
            )
        if right_balance_label:
            _set_text_shape_text(
                right_balance_label,
                f"戶口餘額${gf_data['account_balance_at_target_year']}"
            )
        _apply_withdrawal_highlights(slide_10, withdrawal_picture, gf_data)
    else:
        _delete_slide(prs, 9)

    prs.save(output_path)


# ── Background task workers ──────────────────────────────────

def _run_standard_task(task_id, all_base64_images, prompt, language):
    """Standard (non-hybrid) PPT generation — runs in background thread."""
    try:
        vault_resp = requests.post(
            PPT_AI_ENDPOINT,
            json={
                "prompt": prompt,
                "language": language,
                "images": all_base64_images
            },
            timeout=180
        )

        if vault_resp.status_code == 503:
            raise ValueError('尚未設定 API Key，請聯絡管理員。')
        if vault_resp.status_code >= 400:
            raise ValueError(_format_vault_error(vault_resp, '1006 AI 生成失敗，請檢查 API Key、JSON schema 或素材格式。'))

        vault_data = vault_resp.json()
        ai_content = vault_data.get('content', '')
        try:
            json_data = json.loads(ai_content)
        except json.JSONDecodeError as e:
            raise ValueError(f'AI JSON 格式錯誤：{e}')

        output_filename = f"generated_presentation_{uuid.uuid4()}.pptx"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
        create_pptx_from_json(json_data, output_path, hybrid_mode=False)

        with TASKS_LOCK:
            TASKS[task_id]['status'] = 'done'
            TASKS[task_id]['result'] = {'download_url': f'/api/download/{output_filename}'}

    except Exception as e:
        print(f"[Task {task_id}] Standard task error: {e}")
        with TASKS_LOCK:
            TASKS[task_id]['status'] = 'error'
            TASKS[task_id]['result'] = {'error': str(e) or 'AI 生成失敗，請稍後重試。'}


def _run_hybrid_task(task_id, all_base64_images, prompt, language, aspect_ratio):
    """Hybrid (AI background) PPT generation — runs in background thread."""
    try:
        vault_resp = requests.post(
            PPT_GENERATE_ENDPOINT,
            json={
                "prompt": prompt,
                "language": language,
                "images": all_base64_images,
                "aspect_ratio": aspect_ratio
            },
            timeout=300
        )

        if vault_resp.status_code == 503:
            raise ValueError('尚未設定 API Key，請聯絡管理員。')
        if vault_resp.status_code >= 400:
            raise ValueError(_format_vault_error(vault_resp, '1006 AI 背景簡報生成失敗，請檢查 API Key、JSON schema 或素材格式。'))

        vault_data = vault_resp.json()
        ai_content = vault_data.get('content', '')
        try:
            json_data = json.loads(ai_content)
        except json.JSONDecodeError as e:
            raise ValueError(f'AI JSON 格式錯誤：{e}')

        output_filename = f"hybrid_presentation_{uuid.uuid4()}.pptx"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
        create_pptx_from_json(json_data, output_path, hybrid_mode=True)

        with TASKS_LOCK:
            TASKS[task_id]['status'] = 'done'
            TASKS[task_id]['result'] = {'download_url': f'/api/download/{output_filename}'}

    except Exception as e:
        print(f"[Task {task_id}] Hybrid task error: {e}")
        with TASKS_LOCK:
            TASKS[task_id]['status'] = 'error'
            TASKS[task_id]['result'] = {'error': str(e) or 'AI 簡報生成失敗，請稍後重試。'}


def _run_gf_task(task_id, pdf_path, client_name, agent_name, withdrawal_mode):
    """GF template finalizer — runs in background thread and does not call AI."""
    crop_paths = {}
    try:
        gf_data = _parse_gf_proposal(pdf_path, client_name, agent_name, withdrawal_mode)
        crop_paths = _render_gf_crops(pdf_path, gf_data, task_id)
        if 'proposal_summary' not in crop_paths:
            raise ValueError('建議書摘要截圖裁切失敗。')
        if 'summary_table' not in crop_paths:
            raise ValueError('基本計劃摘要截圖裁切失敗。')
        if gf_data.get('has_withdrawal') and 'withdrawal_surrender_table' not in crop_paths:
            raise ValueError('提款後退保發還金額截圖裁切失敗。')

        output_filename = f"gf_finalized_{uuid.uuid4()}.pptx"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
        _finalize_gf_pptx(gf_data, crop_paths, output_path)
        pdf_path = _convert_pptx_to_pdf_or_raise(output_path)
        pdf_filename = os.path.basename(pdf_path)

        with TASKS_LOCK:
            TASKS[task_id]['status'] = 'done'
            TASKS[task_id]['result'] = {
                'download_url': f'/api/download/{output_filename}',
                'pptx_download_url': f'/api/download/{output_filename}',
                'pdf_download_url': f'/api/download/{pdf_filename}',
                'metadata': {
                    'plan': 'GF',
                    'has_withdrawal': gf_data.get('has_withdrawal'),
                    'growth_15y': gf_data.get('growth_15y'),
                    'growth_25y': gf_data.get('growth_25y'),
                    'withdrawal_target_year': gf_data.get('withdrawal_target_year'),
                }
            }

    except Exception as e:
        print(f"[Task {task_id}] GF task error: {e}")
        with TASKS_LOCK:
            TASKS[task_id]['status'] = 'error'
            TASKS[task_id]['result'] = {'error': str(e) or 'GF PPT 生成失敗，請檢查建議書格式。'}
    finally:
        for path in [pdf_path, *crop_paths.values()]:
            try:
                if path and os.path.exists(path):
                    os.remove(path)
            except Exception:
                pass


# ── Routes ────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/tools/gf')
def gf_tool():
    return render_template('gf.html')


@app.route('/gf')
def gf_shortcut():
    return redirect('/tools/gf')


@app.route('/tools/gf-model')
def gf_model_tool():
    if not os.path.exists(GF_MODEL_PAGE_PATH):
        return (
            '<!doctype html><html lang="zh-TW"><head><meta charset="utf-8">'
            '<title>GF 提款規律研究</title></head><body>'
            '<h1>GF 提款規律研究工具未安裝</h1>'
            '<p>請先把 gf_withdrawal_model.html 放到 AIAtools/research/。</p>'
            '</body></html>',
            404,
            {'Content-Type': 'text/html; charset=utf-8'},
        )
    return send_file(GF_MODEL_PAGE_PATH)


@app.route('/api/ping')
def ping():
    """Session keepalive — called by frontend every 60s to refresh session TTL."""
    sid = request.cookies.get('auth_sid')
    if sid and sid in AUTH_SESSIONS:
        AUTH_SESSIONS[sid]['expiry'] = datetime.now() + timedelta(seconds=SESSION_TTL_SECONDS)
    return jsonify({'ok': True})


@app.route('/api/task/<task_id>')
def get_task_status(task_id):
    """Poll for async task result."""
    _cleanup_old_tasks()
    with TASKS_LOCK:
        task = TASKS.get(task_id)
    if not task:
        return jsonify({'status': 'not_found'}), 404
    resp = {'status': task['status']}
    if task['result']:
        resp.update(task['result'])
    return jsonify(resp)


@app.route('/api/gf/generate', methods=['POST'])
def generate_gf_ppt():
    """
    GF Finalizer: fill the fixed GF template from a proposal PDF.
    """
    client_name = request.form.get('client_name', '').strip()
    agent_name = request.form.get('agent_name', '').strip()
    withdrawal_mode = request.form.get('withdrawal_mode', 'auto').strip().lower()
    proposal = request.files.get('proposal_pdf') or request.files.get('file')

    if not client_name:
        return jsonify({'error': '請輸入客人名。'}), 400
    if not agent_name:
        return jsonify({'error': '請輸入 Agent 名。'}), 400
    if withdrawal_mode not in {'auto', 'yes', 'no'}:
        withdrawal_mode = 'auto'
    if not proposal or not proposal.filename:
        return jsonify({'error': '請上傳 GF 建議書 PDF。'}), 400
    if _filename_extension(proposal.filename) != 'pdf':
        return jsonify({'error': 'GF Finalizer 只接受 PDF 建議書。'}), 400

    task_id = str(uuid.uuid4())
    safe_filename = f"gf_{task_id}.pdf"
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)

    try:
        proposal.save(pdf_path)
    except Exception:
        return jsonify({'error': 'PDF 上傳失敗，請重新提交。'}), 400

    with TASKS_LOCK:
        TASKS[task_id] = {'status': 'processing', 'result': None, 'created': datetime.now()}

    t = threading.Thread(
        target=_run_gf_task,
        args=(task_id, pdf_path, client_name, agent_name, withdrawal_mode),
        daemon=True
    )
    t.start()

    return jsonify({'task_id': task_id, 'status': 'processing'})


@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    """
    Standard mode: prompt is required; files are optional source material.
    """
    files = request.files.getlist('files') if 'files' in request.files else []
    prompt = request.form.get('prompt', '').strip()
    language = request.form.get('language', 'Traditional Chinese')

    if not prompt:
        return jsonify({'error': '請輸入核心提示詞，說明你想做成什麼 PPT。'}), 400

    try:
        all_base64_images, processed_any_file = _collect_uploaded_images(files)
    except ValueError as e:
        return jsonify({'error': str(e)}), 400
    if processed_any_file and not all_base64_images:
        return jsonify({'error': '素材解析失敗。請確認檔案為有效 PDF、PPT、PPTX、PNG 或 JPG。'}), 400

    task_id = str(uuid.uuid4())
    with TASKS_LOCK:
        TASKS[task_id] = {'status': 'processing', 'result': None, 'created': datetime.now()}

    t = threading.Thread(
        target=_run_standard_task,
        args=(task_id, all_base64_images, prompt, language),
        daemon=True
    )
    t.start()

    return jsonify({'task_id': task_id, 'status': 'processing'})


@app.route('/api/generate-ppt/hybrid', methods=['POST'])
def generate_ppt_hybrid():
    """
    Hybrid mode: prompt is required; files are optional source material.
    """
    files = request.files.getlist('files') if 'files' in request.files else []
    prompt = request.form.get('prompt', '').strip()
    language = request.form.get('language', 'Traditional Chinese')
    aspect_ratio = request.form.get('aspect_ratio', '16:9')

    if aspect_ratio not in ASPECT_RATIO_MAP:
        aspect_ratio = '16:9'

    if not prompt:
        return jsonify({'error': '請輸入核心提示詞，說明你想做成什麼 PPT。'}), 400

    try:
        all_base64_images, processed_any_file = _collect_uploaded_images(files)
    except ValueError as e:
        return jsonify({'error': str(e)}), 400
    if processed_any_file and not all_base64_images:
        return jsonify({'error': '素材解析失敗。請確認檔案為有效 PDF、PPT、PPTX、PNG 或 JPG。'}), 400

    task_id = str(uuid.uuid4())
    with TASKS_LOCK:
        TASKS[task_id] = {'status': 'processing', 'result': None, 'created': datetime.now()}

    t = threading.Thread(
        target=_run_hybrid_task,
        args=(task_id, all_base64_images, prompt, language, aspect_ratio),
        daemon=True
    )
    t.start()

    return jsonify({'task_id': task_id, 'status': 'processing'})


@app.route('/api/download/<filename>')
def download_file(filename):
    filename = os.path.basename(filename)
    filepath = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    if filename.lower().endswith('.pdf') and not os.path.exists(filepath):
        pptx_filename = filename[:-4] + '.pptx'
        pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], pptx_filename)
        if os.path.exists(pptx_path):
            try:
                filepath = _convert_pptx_to_pdf_or_raise(pptx_path)
            except Exception as e:
                return jsonify({'error': str(e) or 'PDF conversion failed'}), 500
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({'error': 'File not found'}), 404


@app.route('/health')
def health():
    return jsonify({'status': 'ok', 'service': '1008-ppt-generator'}), 200


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5008, debug=True)
